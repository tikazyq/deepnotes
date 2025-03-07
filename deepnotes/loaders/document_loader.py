import os
from typing import Optional

from deepnotes.loaders.base_loader import BaseLoader
from deepnotes.models.loader_models import (
    DocumentChunk,
    DocumentLoadedData,
    FileMetadata,
    ProcessedDocument,
)
from deepnotes.models.source_models import (
    SourceConfig,
    SourceType,
)

# PDF processing: PyPDF2
try:
    import PyPDF2
    from PyPDF2 import PdfReader
except ImportError:
    PyPDF2 = None
    PdfReader = None

# Word processing: python-docx
try:
    from docx import Document
except ImportError:
    Document = None

# PPTX processing: python-pptx
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

# Excel processing: pandas
try:
    import pandas as pd
except ImportError:
    pd = None

# Text splitting from llama-index
try:
    from llama_index.core.node_parser import TokenTextSplitter
except ImportError:
    TokenTextSplitter = None


class DocumentLoader(BaseLoader):
    def __init__(self, config: SourceConfig):
        super().__init__(config)

    @classmethod
    def get_source_type(cls) -> SourceType:
        return SourceType.DOCUMENT

    def validate_config(self) -> bool:
        return "path" in self.config.connection

    @staticmethod
    def _extract_pdf(file_path: str) -> str:
        """Extract text from PDF using PyPDF2"""
        text = []
        with open(file_path, "rb") as f:
            reader = PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text.append(page_text)
        return "\n".join(text)

    @staticmethod
    def _extract_docx(file_path: str) -> str:
        """Extract text from Word documents"""
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])

    @staticmethod
    def _extract_pptx(file_path: str) -> str:
        """Extract text from PowerPoint presentations"""
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    @staticmethod
    def _extract_excel(file_path: str) -> str:
        """Extract text from Excel files using pandas DataFrames"""
        xls = pd.ExcelFile(file_path)
        text = []
        for sheet_name in xls.sheet_names:
            df = xls.parse(sheet_name)
            text.append(f"=== Sheet: {sheet_name} ===\n{df.to_string(index=False)}")
        return "\n\n".join(text)

    def process(self, *, target_path: Optional[str]=None) -> DocumentLoadedData:
        """Process a file or directory"""
        if not target_path:
            target_path = self.config.connection["path"]

        if not os.path.exists(target_path):
            raise FileNotFoundError(f"Path not found: {target_path}")

        loaded_docs = []
        file_types = set()

        if os.path.isdir(target_path):
            for root, _, files in os.walk(target_path):
                for file in files:
                    file_path = os.path.join(root, file)
                    loaded_docs.append(self._process_single_file(file_path))
                    file_types.add(os.path.splitext(file)[1].lower())
        else:
            loaded_docs.append(self._process_single_file(target_path))
            file_types.add(os.path.splitext(target_path)[1].lower())

        return DocumentLoadedData(
            source_type=SourceType.DOCUMENT,
            global_metadata={
                "source_path": target_path,
                "file_types": list(file_types),
                "total_chunks": sum(len(d.chunks) for d in loaded_docs),
            },
            documents=loaded_docs,
        )

    def _process_single_file(self, file_path: str) -> ProcessedDocument:
        """Process individual file and return structured data"""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Document not found: {file_path}")

        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            if not PyPDF2:
                raise ImportError("PyPDF2 required for PDF processing")
            content = self._extract_pdf(file_path)
        elif ext == ".docx":
            if not Document:
                raise ImportError("python-docx required for Word processing")
            content = self._extract_docx(file_path)
        elif ext == ".pptx":
            if not Presentation:
                raise ImportError("python-pptx required for PPTX processing")
            content = self._extract_pptx(file_path)
        elif ext in (".xlsx", ".xls"):
            if not pd:
                raise ImportError("pandas required for Excel processing")
            content = self._extract_excel(file_path)
        elif ext in (".txt", ".md"):
            with open(file_path, 'r') as f:
                content = f.read()
        else:
            # Attempt to read as plain text if content is decodable
            try:
                # Verify text validity by attempting UTF-8 decode
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
            except UnicodeDecodeError:
                raise ValueError(f"File appears to be binary: {file_path}")
            except Exception as e:
                raise ValueError(f"Error reading file {file_path}: {str(e)}")

        if not TokenTextSplitter:
            raise ImportError("llama-index required for text splitting")
        splitter = TokenTextSplitter(
            chunk_size=self.config.options.get("chunk_size", 2000),
            chunk_overlap=self.config.options.get("chunk_overlap", 200),
        )
        chunks = splitter.split_text(content)

        return ProcessedDocument(
            metadata=FileMetadata(
                file_path=file_path,
                file_size=os.path.getsize(file_path),
                file_type=ext,
                created_at=str(os.path.getctime(file_path)),
                updated_at=str(os.path.getmtime(file_path)),
            ),
            chunks=[
                DocumentChunk(text=chunk, index=idx) for idx, chunk in enumerate(chunks)
            ],
            raw_content=content,
        )
