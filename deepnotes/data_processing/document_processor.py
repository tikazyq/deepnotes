import os
from typing import Optional

from deepnotes.data_processing import BaseDataProcessor, DataSourceConfig
from deepnotes.models import IntermediateDataModel

# PDF processing: PyPDF2
try:
    import PyPDF2
    from PyPDF2 import PdfReader
except ImportError:
    PyPDF2 = None
    PdfReader = None

# Word processing: python-docx
try:
    from docx import Document
except ImportError:
    Document = None

# PPTX processing: python-pptx
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

# Excel processing: pandas
try:
    import pandas as pd
except ImportError:
    pd = None

# Text splitting from llama-index
try:
    from llama_index.core.node_parser import TokenTextSplitter
except ImportError:
    TokenTextSplitter = None


class DocumentProcessorConfig(DataSourceConfig):
    chunk_size: int = 1000
    chunk_overlap: int = 200


class DocumentProcessor(BaseDataProcessor):
    def __init__(self, config: DocumentProcessorConfig):
        super().__init__(config)
        self.raw_text: Optional[str] = None

    def _extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF using PyPDF2"""
        text = []
        with open(file_path, "rb") as f:
            reader = PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text.append(page_text)
        return "\n".join(text)

    def _extract_docx(self, file_path: str) -> str:
        """Extract text from Word documents"""
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])

    def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint presentations"""
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    def _extract_excel(self, file_path: str) -> str:
        """Extract text from Excel files using pandas DataFrames"""
        xls = pd.ExcelFile(file_path)
        text = []
        for sheet_name in xls.sheet_names:
            df = xls.parse(sheet_name)
            text.append(f"=== Sheet: {sheet_name} ===\n{df.to_string(index=False)}")
        return "\n\n".join(text)

    def extract(self) -> str:
        """Extract raw text from document based on file type"""
        file_path = self.config.source_path
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Document not found: {file_path}")

        ext = os.path.splitext(file_path)[1].lower()

        if ext == ".pdf":
            if not PyPDF2:
                raise ImportError("PyPDF2 required for PDF processing")
            return self._extract_pdf(file_path)

        elif ext == ".docx":
            if not Document:
                raise ImportError("python-docx required for Word processing")
            return self._extract_docx(file_path)

        elif ext == ".pptx":
            if not Presentation:
                raise ImportError("python-pptx required for PPTX processing")
            return self._extract_pptx(file_path)

        elif ext in (".xlsx", ".xls"):
            if not pd:
                raise ImportError("pandas required for Excel processing")
            return self._extract_excel(file_path)

        raise ValueError(f"Unsupported file type: {ext}")

    def process(self) -> IntermediateDataModel:
        """Process document into chunks using text splitting"""
        if not TokenTextSplitter:
            raise ImportError("llama-index required for text splitting")

        raw_text = self.extract()
        splitter = TokenTextSplitter(
            chunk_size=self.config.chunk_size, chunk_overlap=self.config.chunk_overlap
        )
        chunks = splitter.split_text(raw_text)

        return IntermediateDataModel(
            entities=[],
            relationships=[],
            raw_data=[{"text": chunk} for chunk in chunks],
            notes=[f"Extracted {len(chunks)} chunks from document"],
            iteration_history=[],
            metadata=[],
        )
        # return {
        #     "document_type": os.path.splitext(self.config.source_path)[1].lower(),
        #     "source_path": self.config.source_path,
        #     "chunks": chunks,
        #     "chunk_size": self.config.chunk_size,
        #     "chunk_overlap": self.config.chunk_overlap,
        #     "total_chunks": len(chunks),
        # }
